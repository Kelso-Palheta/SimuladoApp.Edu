"""
gerador_pptx.py — Gerador de versão PPTX editável.
Requer: pip install python-pptx
"""
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("ERRO: Biblioteca 'python-pptx' não está instalada. Execute: pip install python-pptx")
    import sys
    sys.exit(1)

def gerar_pptx(dados_aula, nome_arquivo="aula_gerada.pptx"):
    prs = Presentation()

    # Slide 1: Capa
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = dados_aula.get("titulo", "Título da Aula")
    subtitle.text = f"{dados_aula.get('subtitulo', '')}\n\n{dados_aula.get('materia', '')} - {dados_aula.get('serie', '')}"

    # Fundo escuro na capa
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(11, 12, 16) # #0B0C10
    
    # Slides de Conteúdo
    for conteudo in dados_aula.get("conteudo", []):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Fundo
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(11, 12, 16)

        title = slide.shapes.title
        body = slide.placeholders[1]
        
        title.text = conteudo.get("titulo", "")
        # Ajusta a cor do titulo
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 240, 255) # #00F0FF Cyberpunk
        
        tf = body.text_frame
        tf.text = conteudo.get("texto", "")
        for p in tf.paragraphs:
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(18)

        # Callouts simples no PPTX
        callouts = conteudo.get("callouts", [])
        if callouts:
            p = tf.add_paragraph()
            p.text = "\nDESTAQUES:"
            p.font.bold = True
            p.font.color.rgb = RGBColor(229, 9, 20) # #E50914 Netflix Red
            
            for c in callouts:
                p = tf.add_paragraph()
                p.text = f"[{c.get('label', '')}] {c.get('texto', '')}"
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(200, 200, 200)

    prs.save(nome_arquivo)
    print(f"✓ PPTX Base gerado com sucesso: {nome_arquivo}")
    print("DICA: Abra no PowerPoint, selecione todos os slides e aplique a Transição 'Transformar/Morph' para criar o efeito contínuo!")
